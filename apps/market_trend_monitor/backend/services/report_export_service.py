"""レポート二進制エクスポートサービス.

Market Trend Monitor のレポートを企業向けテンプレートで
PDF / PPTX 形式へ出力します。
"""

from __future__ import annotations

import io
import logging
import re
from datetime import datetime
from html import escape
from typing import Any


class ReportExportService:
    """レポートのPDF/PPTX生成サービス."""

    def __init__(self) -> None:
        self._logger = logging.getLogger(self.__class__.__name__)
        self._has_reportlab = self._check_reportlab()
        self._has_python_pptx = self._check_python_pptx()

    def _check_reportlab(self) -> bool:
        """ReportLab 利用可否を確認."""
        try:
            from reportlab.lib.pagesizes import A4  # noqa: F401

            return True
        except ImportError:
            self._logger.warning("ReportLab が未インストールのため PDF 出力は利用不可")
            return False

    def _check_python_pptx(self) -> bool:
        """python-pptx 利用可否を確認."""
        try:
            from pptx import Presentation  # noqa: F401

            return True
        except ImportError:
            self._logger.warning("python-pptx が未インストールのため PPTX 出力は利用不可")
            return False

    def export_pdf(self, report: dict[str, Any]) -> bytes:
        """PDFバイナリを生成."""
        if not self._has_reportlab:
            raise RuntimeError("ReportLab が未インストールのため PDF 出力できません")

        try:
            return self._build_pdf(report)
        except Exception as exc:
            self._logger.error("PDF生成に失敗: %s", exc, exc_info=True)
            raise RuntimeError(f"PDF生成に失敗しました: {exc}") from exc

    def export_pptx(self, report: dict[str, Any]) -> bytes:
        """PPTXバイナリを生成."""
        if not self._has_python_pptx:
            raise RuntimeError("python-pptx が未インストールのため PPTX 出力できません")

        try:
            return self._build_pptx(report)
        except Exception as exc:
            self._logger.error("PPTX生成に失敗: %s", exc, exc_info=True)
            raise RuntimeError(f"PPTX生成に失敗しました: {exc}") from exc

    def _build_pdf(self, report: dict[str, Any]) -> bytes:
        """企業テンプレートPDFを生成."""
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import cm
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.cidfonts import UnicodeCIDFont
        from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
        from reportlab.graphics.charts.barcharts import VerticalBarChart
        from reportlab.graphics.shapes import Drawing

        font_name = "Helvetica"
        try:
            pdfmetrics.registerFont(UnicodeCIDFont("HeiseiMin-W3"))
            font_name = "HeiseiMin-W3"
        except Exception:
            self._logger.warning("CJKフォント登録に失敗。PDF表示品質が低下する可能性があります")

        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            "ReportTitle",
            parent=styles["Title"],
            fontName=font_name,
            fontSize=24,
            spaceAfter=14,
        )
        h1_style = ParagraphStyle(
            "ReportHeading1",
            parent=styles["Heading2"],
            fontName=font_name,
            fontSize=16,
            spaceAfter=10,
            spaceBefore=8,
        )
        normal_style = ParagraphStyle(
            "ReportBody",
            parent=styles["Normal"],
            fontName=font_name,
            fontSize=10,
            leading=14,
            spaceAfter=6,
        )

        title = str(report.get("title", "市場動向レポート"))
        report_id = str(report.get("id", "-"))
        created_at = self._safe_datetime(str(report.get("created_at", "")))
        period_start = self._safe_datetime(str(report.get("period_start", "")))
        period_end = self._safe_datetime(str(report.get("period_end", "")))
        sections = self._sections(report)
        trends = self._trends(report)

        avg_score = sum(self._as_float(t.get("score")) for t in trends) / max(len(trends), 1)
        total_articles = sum(self._as_int(t.get("articles_count")) for t in trends)

        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4, topMargin=2 * cm, bottomMargin=1.7 * cm)
        elements: list[Any] = []

        # 1. 表紙
        elements.append(Paragraph(escape(title), title_style))
        elements.append(Paragraph("企業向け 市場動向分析レポート", normal_style))
        elements.append(Spacer(1, 0.8 * cm))
        elements.append(Paragraph(f"レポートID: {escape(report_id)}", normal_style))
        elements.append(Paragraph(f"作成日時: {escape(created_at)}", normal_style))
        elements.append(
            Paragraph(
                f"対象期間: {escape(period_start)} - {escape(period_end)}",
                normal_style,
            )
        )
        elements.append(Spacer(1, 1.2 * cm))
        elements.append(
            Paragraph(
                "本資料は市場監視ダッシュボードの分析結果を正式配布用に再構成したものです。",
                normal_style,
            )
        )
        elements.append(PageBreak())

        # 2. 目次
        toc_items = [
            "1. 表紙",
            "2. 目次",
            "3. KPIサマリー",
            "4. トレンドチャート",
            "5. 詳細分析セクション",
            "6. 結論と次アクション",
        ]
        elements.append(Paragraph("目次", h1_style))
        for item in toc_items:
            elements.append(Paragraph(item, normal_style))
        elements.append(PageBreak())

        # 3. KPIサマリー
        elements.append(Paragraph("KPIサマリー", h1_style))
        kpi_table = Table(
            [
                ["指標", "値"],
                ["トレンド件数", str(len(trends))],
                ["総記事件数", str(total_articles)],
                ["平均スコア", f"{avg_score:.2f}"],
                ["セクション件数", str(len(sections))],
            ],
            colWidths=[6.5 * cm, 8.5 * cm],
        )
        kpi_table.setStyle(
            TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#E2E8F0")),
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#CBD5E1")),
                    ("FONTNAME", (0, 0), (-1, -1), font_name),
                    ("FONTSIZE", (0, 0), (-1, -1), 10),
                    ("ALIGN", (0, 0), (-1, -1), "LEFT"),
                ]
            )
        )
        elements.append(kpi_table)
        elements.append(Spacer(1, 0.6 * cm))
        elements.append(Paragraph("注: NEW は前期間比較の基準データ不足を示します。", normal_style))
        elements.append(PageBreak())

        # 4. チャート
        elements.append(Paragraph("トレンドチャート", h1_style))
        chart_topics = [self._truncate(str(item.get("topic", "-")), 16) for item in trends[:8]]
        chart_scores = [self._as_float(item.get("score")) for item in trends[:8]]

        if chart_topics:
            drawing = Drawing(460, 240)
            bar = VerticalBarChart()
            bar.x = 30
            bar.y = 45
            bar.height = 160
            bar.width = 390
            bar.data = [chart_scores]
            bar.categoryAxis.categoryNames = chart_topics
            bar.valueAxis.valueMin = 0
            bar.valueAxis.valueMax = max(max(chart_scores), 1.0)
            bar.valueAxis.valueStep = max(bar.valueAxis.valueMax / 5, 0.2)
            bar.bars[0].fillColor = colors.HexColor("#2563EB")
            drawing.add(bar)
            elements.append(drawing)

        trend_table_rows: list[list[str]] = [["トピック", "スコア", "記事数", "成長"]]
        for trend in trends[:10]:
            trend_table_rows.append(
                [
                    str(trend.get("topic", "-")),
                    f"{self._as_float(trend.get('score')):.2f}",
                    str(self._as_int(trend.get("articles_count"))),
                    self._growth_label(trend),
                ]
            )

        trend_table = Table(trend_table_rows, colWidths=[6.8 * cm, 2.5 * cm, 2.5 * cm, 3.2 * cm])
        trend_table.setStyle(
            TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#DBEAFE")),
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#BFDBFE")),
                    ("FONTNAME", (0, 0), (-1, -1), font_name),
                    ("FONTSIZE", (0, 0), (-1, -1), 9),
                ]
            )
        )
        elements.append(trend_table)
        elements.append(PageBreak())

        # 5. 詳細分析
        elements.append(Paragraph("詳細分析セクション", h1_style))
        if not sections:
            elements.append(Paragraph("出力可能なセクションがありません。", normal_style))
        for index, section in enumerate(sections, start=1):
            section_title = str(section.get("title", f"Section {index}"))
            section_content = self._sanitize_markdown(str(section.get("content", "")))
            elements.append(Paragraph(f"{index}. {escape(section_title)}", h1_style))
            elements.append(Paragraph(escape(section_content), normal_style))
            if index % 2 == 0 and index != len(sections):
                elements.append(PageBreak())

        elements.append(PageBreak())

        # 6. 結論
        elements.append(Paragraph("結論と次アクション", h1_style))
        top_trend = trends[0] if trends else {}
        top_topic = str(top_trend.get("topic", "主要トピックなし"))
        top_score = self._as_float(top_trend.get("score"))
        top_growth = self._growth_label(top_trend)

        conclusion_lines = [
            f"最注目トピック: {top_topic} (スコア {top_score:.2f}, 成長 {top_growth})",
            "推奨アクション1: NEW 判定トピックは母数不足を確認してから意思決定する。",
            "推奨アクション2: 上位3トピックの追加証拠を次周期で重点収集する。",
            "推奨アクション3: レポート配布前に数値説明（成長式・基準期間）を添付する。",
        ]
        for line in conclusion_lines:
            elements.append(Paragraph(escape(line), normal_style))

        doc.build(elements)
        return buffer.getvalue()

    def _build_pptx(self, report: dict[str, Any]) -> bytes:
        """企業テンプレートPPTXを生成."""
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.dml.color import RGBColor
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches, Pt

        title = str(report.get("title", "市場動向レポート"))
        report_id = str(report.get("id", "-"))
        created_at = self._safe_datetime(str(report.get("created_at", "")))
        period_start = self._safe_datetime(str(report.get("period_start", "")))
        period_end = self._safe_datetime(str(report.get("period_end", "")))
        sections = self._sections(report)
        trends = self._trends(report)

        presentation = Presentation()

        # 1. 表紙
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        slide.shapes.title.text = title
        subtitle = slide.placeholders[1]
        subtitle.text = (
            "企業向け 市場動向分析レポート\n"
            f"レポートID: {report_id}\n"
            f"作成日時: {created_at}\n"
            f"対象期間: {period_start} - {period_end}"
        )

        # 2. 目次
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "目次"
        content_frame = slide.placeholders[1].text_frame
        content_frame.clear()
        toc_items = [
            "1. KPIサマリー",
            "2. トレンドチャート",
            "3. 成長率チャート",
            "4. 詳細分析",
            "5. 結論と次アクション",
        ]
        for item in toc_items:
            paragraph = content_frame.add_paragraph()
            paragraph.text = item
            paragraph.level = 0

        # 3. KPIサマリー
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "KPIサマリー"
        textbox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.2), Inches(3.5))
        tf = textbox.text_frame
        tf.word_wrap = True
        avg_score = sum(self._as_float(t.get("score")) for t in trends) / max(len(trends), 1)
        total_articles = sum(self._as_int(t.get("articles_count")) for t in trends)
        lines = [
            f"トレンド件数: {len(trends)}",
            f"総記事件数: {total_articles}",
            f"平均スコア: {avg_score:.2f}",
            f"分析セクション件数: {len(sections)}",
            "注: NEW は前期間データがない新規検知",
        ]
        tf.text = lines[0]
        for line in lines[1:]:
            p = tf.add_paragraph()
            p.text = line

        # 4. トレンドスコア棒グラフ
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "トレンドスコア"
        score_data = CategoryChartData()
        score_data.categories = [self._truncate(str(t.get("topic", "-")), 18) for t in trends[:8]]
        score_data.add_series("スコア", [self._as_float(t.get("score")) for t in trends[:8]])
        score_chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.7),
            Inches(1.5),
            Inches(8.6),
            Inches(4.3),
            score_data,
        ).chart
        score_chart.has_legend = False
        score_chart.value_axis.maximum_scale = max(1.0, score_chart.value_axis.maximum_scale or 1.0)

        # 5. 成長率折れ線グラフ
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "成長率"
        growth_data = CategoryChartData()
        growth_data.categories = [self._truncate(str(t.get("topic", "-")), 18) for t in trends[:8]]
        growth_data.add_series(
            "成長率(%)",
            [self._growth_value_percent(t) for t in trends[:8]],
        )
        growth_chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE_MARKERS,
            Inches(0.7),
            Inches(1.5),
            Inches(8.6),
            Inches(4.3),
            growth_data,
        ).chart
        growth_chart.has_legend = False

        # 6. 詳細分析と結論
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "結論と次アクション"
        tf = slide.placeholders[1].text_frame
        tf.clear()

        top_trend = trends[0] if trends else {}
        top_topic = str(top_trend.get("topic", "主要トピックなし"))
        top_score = self._as_float(top_trend.get("score"))
        top_growth = self._growth_label(top_trend)

        conclusion_lines = [
            f"最注目トピック: {top_topic} (スコア {top_score:.2f}, 成長 {top_growth})",
            "NEW 判定トピックは母数確認後に優先順位を決定する",
            "上位3トピックの追加証拠収集を次周期計画へ反映する",
        ]
        if sections:
            section_titles = ", ".join(
                self._truncate(str(s.get("title", "")), 24) for s in sections[:3]
            )
            conclusion_lines.append(f"主要分析セクション: {section_titles}")

        tf.text = conclusion_lines[0]
        first_run = tf.paragraphs[0].runs[0]
        first_run.font.bold = True
        first_run.font.size = Pt(20)
        first_run.font.color.rgb = RGBColor(30, 64, 175)
        for line in conclusion_lines[1:]:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(16)

        output = io.BytesIO()
        presentation.save(output)
        return output.getvalue()

    def _sections(self, report: dict[str, Any]) -> list[dict[str, Any]]:
        """セクション一覧を正規化."""
        sections = report.get("sections", [])
        if not isinstance(sections, list):
            return []
        return [item for item in sections if isinstance(item, dict)]

    def _trends(self, report: dict[str, Any]) -> list[dict[str, Any]]:
        """トレンド一覧を正規化."""
        trends = report.get("trends", [])
        if not isinstance(trends, list):
            return []
        return [item for item in trends if isinstance(item, dict)]

    def _safe_datetime(self, raw: str) -> str:
        """日時文字列を安全に表示形式へ変換."""
        if not raw:
            return "-"

        try:
            normalized = raw.replace("Z", "+00:00")
            value = datetime.fromisoformat(normalized)
            return value.strftime("%Y-%m-%d %H:%M")
        except ValueError:
            return raw

    def _as_float(self, value: Any) -> float:
        """値をfloatへ安全変換."""
        try:
            return float(value)
        except (TypeError, ValueError):
            return 0.0

    def _as_int(self, value: Any) -> int:
        """値をintへ安全変換."""
        try:
            return int(value)
        except (TypeError, ValueError):
            return 0

    def _truncate(self, text: str, size: int) -> str:
        """文字列を指定長で省略."""
        if len(text) <= size:
            return text
        return f"{text[:size]}..."

    def _sanitize_markdown(self, text: str) -> str:
        """Markdownをプレーンテキストへ簡易変換."""
        normalized = text.replace("\r", "")
        normalized = re.sub(r"^#{1,6}\s+", "", normalized, flags=re.MULTILINE)
        normalized = re.sub(r"^[-*]\s+", "", normalized, flags=re.MULTILINE)
        normalized = re.sub(r"^\d+\.\s+", "", normalized, flags=re.MULTILINE)
        normalized = re.sub(r"\*\*(.*?)\*\*", r"\1", normalized)
        normalized = re.sub(r"`(.*?)`", r"\1", normalized)
        normalized = re.sub(r"\n{3,}", "\n\n", normalized)
        return normalized.strip() or "(本文なし)"

    def _growth_label(self, trend: dict[str, Any]) -> str:
        """成長率表示ラベルを返す."""
        metadata = trend.get("metadata", {})
        growth_state = metadata.get("growth_state") if isinstance(metadata, dict) else None
        if growth_state == "new":
            return "NEW"
        if growth_state in {"insufficient_history", "no_signal"}:
            return "N/A"
        return f"{self._as_float(trend.get('growth_rate')) * 100:.1f}%"

    def _growth_value_percent(self, trend: dict[str, Any]) -> float:
        """チャート用成長率(%)."""
        metadata = trend.get("metadata", {})
        growth_state = metadata.get("growth_state") if isinstance(metadata, dict) else None
        if growth_state == "new":
            return 0.0
        return self._as_float(trend.get("growth_rate")) * 100.0
